"""PPTX ingestion adapter — one record per slide text block."""

from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Iterable, Iterator

from agents.ingestion.base import (
    IngestionAdapter,
    RawRecord,
    Source,
    SourceManifest,
    hash_content,
    make_source_record_id,
)

try:
    from pptx import Presentation
    _PPTX_AVAILABLE = True
except ImportError:
    _PPTX_AVAILABLE = False


class PptxAdapter(IngestionAdapter):
    source_type = "docx"
    schema_version = "1.0"
    MIN_LEN = 20

    def discover(self, root: Path) -> Iterable[Source]:
        for path in sorted(root.glob("**/*.pptx")):
            file_hash = hashlib.sha256(path.read_bytes()).hexdigest()
            yield Source(
                path=path,
                source_type=self.source_type,
                relative_path=str(path.relative_to(root)),
                file_hash=file_hash,
            )

    def extract_raw(self, source: Source) -> Iterator[RawRecord]:
        if not _PPTX_AVAILABLE:
            yield _error(source, "python-pptx not installed")
            return
        try:
            prs = Presentation(str(source.path))
        except Exception as exc:
            yield _error(source, str(exc))
            return

        for slide_idx, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text and shape.text.strip():
                    texts.append(shape.text.strip())
            if not texts:
                continue
            combined = "\n".join(texts)
            if len(combined) < self.MIN_LEN:
                continue
            chunk = {
                "text": combined,
                "_slide": slide_idx,
                "_format": "pptx",
                "_extraction_mode": "slide_text",
            }
            offset = f"slide:{slide_idx}"
            ch = hash_content(chunk)
            yield RawRecord(
                source_record_id=make_source_record_id(source.relative_path, offset, ch),
                source_file=source.relative_path,
                source_type=self.source_type,
                source_offset=offset,
                content_hash=ch,
                raw_content=chunk,
            )

    def manifest(self, source: Source) -> SourceManifest:
        records = list(self.extract_raw(source))
        return SourceManifest(
            source_file=source.relative_path,
            source_type=self.source_type,
            file_hash=source.file_hash,
            record_count=len(records),
        )


def _error(source: Source, msg: str) -> RawRecord:
    content = {"_error": msg, "_format": "pptx"}
    ch = hash_content(content)
    offset = "slide:0:error"
    return RawRecord(
        source_record_id=make_source_record_id(source.relative_path, offset, ch),
        source_file=source.relative_path,
        source_type=source.source_type,
        source_offset=offset,
        content_hash=ch,
        raw_content=content,
    )
